from __future__ import annotations

from pathlib import Path
import re
from typing import Iterable

import fitz
from pptx import Presentation


def _normalize_text(text: str) -> str:
    # Keep paragraph breaks (for better chunking) but normalize spaces.
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _chunk_text(text: str, *, chunk_size: int = 1000) -> list[str]:
    """
    Chunk text into <= chunk_size characters, trying to respect paragraph boundaries.
    """
    text = _normalize_text(text)
    if not text:
        return []

    if chunk_size <= 0:
        return [text]

    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p and p.strip()]
    chunks: list[str] = []
    buf: list[str] = []
    buf_len = 0

    def flush() -> None:
        nonlocal buf, buf_len
        if not buf:
            return
        chunk = "\n\n".join(buf).strip()
        if chunk:
            chunks.append(chunk)
        buf = []
        buf_len = 0

    for p in paragraphs:
        # If a single paragraph is too long, hard-split it.
        if len(p) > chunk_size:
            flush()
            start = 0
            while start < len(p):
                chunks.append(p[start : start + chunk_size].strip())
                start += chunk_size
            continue

        # Greedy pack paragraphs into a chunk.
        extra = (2 if buf else 0) + len(p)
        if buf_len + extra > chunk_size:
            flush()
        buf.append(p)
        buf_len += extra

    flush()
    return chunks


def _extract_pdf_segments_with_langchain(file_path: str) -> list[str]:
    """
    Extract PDF text using LangChain's PyPDFLoader, returning per-page segments.
    """
    try:
        # Newer versions typically expose this via langchain_community.
        from langchain_community.document_loaders import PyPDFLoader  # type: ignore
    except Exception as exc:  # noqa: BLE001
        raise ImportError(
            "PyPDFLoader not available. Install `langchain-community` and `pypdf`."
        ) from exc

    loader = PyPDFLoader(file_path)
    docs = loader.load()  # one Document per page in most cases
    segments: list[str] = []
    for doc in docs:
        page_text = getattr(doc, "page_content", "") or ""
        page_text = _normalize_text(page_text)
        if page_text:
            segments.append(page_text)
    return segments


def _extract_pdf_segments_fallback(file_path: str) -> list[str]:
    """
    Fallback extractor (PyMuPDF/fitz), returning per-page segments.
    """
    path = Path(file_path)
    if not path.exists():
        return []

    segments: list[str] = []
    with fitz.open(path) as doc:
        for page in doc:
            text = page.get_text("text") or ""
            text = _normalize_text(text)
            if text:
                segments.append(text)
    return segments


def _walk_ppt_shapes(shape) -> Iterable[str]:
    # Recursively walk grouped shapes and collect text.
    if hasattr(shape, "text") and getattr(shape, "text", None):
        yield shape.text
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            yield from _walk_ppt_shapes(child)


def _extract_ppt_segments(file_path: str) -> list[str]:
    """
    Extract PPT/PPTX text using python-pptx, returning per-slide segments.
    """
    path = Path(file_path)
    if not path.exists():
        return []

    prs = Presentation(str(path))
    segments: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        pieces: list[str] = []
        for shape in slide.shapes:
            for value in _walk_ppt_shapes(shape):
                cleaned = _normalize_text(str(value))
                if cleaned:
                    pieces.append(cleaned)
        if not pieces:
            continue
        slide_text = "\n".join(pieces).strip()
        if slide_text:
            # Treat each slide as a "chapter/segment".
            segments.append(f"[Slide {idx}]\n{slide_text}")
    return segments


def load_document_chunks(
    file_path: str,
    file_type: str,
    *,
    chunk_size: int = 1000,
) -> list[str]:
    """
    Extract (Extract) and segment (Segment) a document into chunks.

    - PDF: LangChain PyPDFLoader (per-page) with a safe fallback to PyMuPDF.
    - PPT/PPTX: python-pptx (per-slide).
    - Chunking: per segment/chapter, and further split every `chunk_size` chars.
    """
    path = Path(file_path)
    if not path.exists():
        return []

    ft = (file_type or "").lower().strip(".")
    segments: list[str] = []
    if ft == "pdf":
        try:
            segments = _extract_pdf_segments_with_langchain(str(path))
        except ImportError:
            segments = _extract_pdf_segments_fallback(str(path))
    elif ft in {"ppt", "pptx"}:
        segments = _extract_ppt_segments(str(path))
    else:
        return []

    chunks: list[str] = []
    for seg in segments:
        chunks.extend(_chunk_text(seg, chunk_size=chunk_size))

    return chunks

