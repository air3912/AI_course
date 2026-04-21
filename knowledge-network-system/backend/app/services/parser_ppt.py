from pathlib import Path

from pptx import Presentation


def _extract_shape_text(shape) -> list[str]:
    values: list[str] = []
    if hasattr(shape, "text") and shape.text:
        values.append(shape.text)
    if hasattr(shape, "shapes"):
        for item in shape.shapes:
            values.extend(_extract_shape_text(item))
    return values

def parse_ppt_text(file_path: str) -> str:
    path = Path(file_path)
    if not path.exists():
        return ""

    prs = Presentation(path)
    chunks: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            chunks.extend(_extract_shape_text(shape))
    return "\n".join(item.strip() for item in chunks if item and item.strip())
