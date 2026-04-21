from pathlib import Path

import fitz
from pptx import Presentation

from app.services.parser_pdf import parse_pdf_text
from app.services.parser_ppt import parse_ppt_text


def test_parse_pdf_text_reads_real_content(tmp_path: Path) -> None:
    pdf_path = tmp_path / "demo.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Knowledge Graph from PDF")
    doc.save(pdf_path)
    doc.close()

    text = parse_pdf_text(str(pdf_path))
    assert "Knowledge Graph" in text


def test_parse_ppt_text_reads_real_content(tmp_path: Path) -> None:
    ppt_path = tmp_path / "demo.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Machine Learning"
    slide.placeholders[1].text = "Knowledge network extraction"
    prs.save(ppt_path)

    text = parse_ppt_text(str(ppt_path))
    assert "Machine Learning" in text
