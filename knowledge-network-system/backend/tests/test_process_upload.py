from io import BytesIO

import fitz
from fastapi.testclient import TestClient
from pptx import Presentation

from app.db.session import init_db
from app.main import app

init_db()
client = TestClient(app)


def _build_pdf_bytes(text: str) -> bytes:
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), text)
    buffer = doc.tobytes()
    doc.close()
    return buffer


def _build_pptx_bytes(title: str, body: str) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = body
    stream = BytesIO()
    prs.save(stream)
    return stream.getvalue()


def test_upload_process_pdf_returns_graph() -> None:
    file_bytes = _build_pdf_bytes("Deep learning model graph relation extraction")
    response = client.post(
        "/api/v1/upload/process",
        files={"file": ("sample.pdf", file_bytes, "application/pdf")},
    )
    assert response.status_code == 200
    data = response.json()
    assert data["parse_status"] == "completed"
    assert data["graph_status"] == "completed"
    assert isinstance(data["nodes"], list)
    assert isinstance(data["edges"], list)
    assert isinstance(data["keywords"], list)
    assert isinstance(data["entities"], list)
    assert data["text_length"] > 0

    graph_response = client.get(f"/api/v1/graph/snapshot/{data['graph_snapshot_id']}")
    assert graph_response.status_code == 200
    graph_payload = graph_response.json()
    assert "stats" in graph_payload


def test_upload_process_pptx_returns_graph() -> None:
    file_bytes = _build_pptx_bytes("NLP", "keyword relation graph network")
    response = client.post(
        "/api/v1/upload/process",
        files={
            "file": (
                "sample.pptx",
                file_bytes,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )
    assert response.status_code == 200
    data = response.json()
    assert data["file_type"] == "pptx"
    assert data["text_length"] > 0
    assert len(data["nodes"]) >= 1
